# =========================================================
# DIVU AI — FINAL PERFECT VERSION
# START TO END FULL CODE
# =========================================================

from flask import Flask, render_template, request, jsonify
from flask_cors import CORS

from huggingface_hub import InferenceClient
from werkzeug.utils import secure_filename

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

import replicate
import requests
import urllib.parse
import random
import os
import time

# =========================================================
# FLASK
# =========================================================

app = Flask(__name__)

CORS(app)

# =========================================================
# TOKENS
# =========================================================

HF_TOKEN = os.getenv("HF_TOKEN")
REPLICATE_API_TOKEN = os.getenv("REPLICATE_API_TOKEN")
# =========================================================
# HUGGINGFACE CHAT MODEL
# =========================================================

client = InferenceClient(

    model="meta-llama/Llama-3.1-8B-Instruct",

    token=HF_TOKEN
)

# =========================================================
# FOLDERS
# =========================================================

UPLOAD_FOLDER = "uploads"

OUTPUT_FOLDER = "static/op"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)

os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# =========================================================
# CHAT MEMORY
# =========================================================

chat_history = []

MAX_HISTORY = 20

# =========================================================
# SYSTEM PROMPT
# =========================================================

SYSTEM_PROMPT = """
You are DIVU AI.

An ultra futuristic premium AI assistant created by Kush.

You are:
- intelligent
- futuristic
- cinematic
- premium
- professional
- creative
- highly helpful

You specialize in:
- AI
- coding
- web development
- app development
- Python
- Android
- startups
- futuristic ideas
- education
- productivity

Rules:
- clean formatting
- modern responses
- beautiful explanations
- complete code only
- beginner friendly
- futuristic tone
- no robotic responses
"""

# =========================================================
# HOME
# =========================================================

@app.route("/")
def home():

    return render_template("index.html")

# =========================================================
# CLEAN RESPONSE
# =========================================================

def clean_reply(reply):

    if not reply:

        return ""

    remove_items = [

        "```python",
        "```html",
        "```css",
        "```javascript",
        "```cpp",
        "```java",
        "```json",
        "```"
    ]

    for item in remove_items:

        reply = reply.replace(item, "")

    return reply.strip()

# =========================================================
# CLEAN IMAGE PROMPT
# =========================================================

def clean_image_prompt(prompt):

    prompt = prompt.lower()

    remove_words = [

        "generate image of",
        "generate image",
        "create image of",
        "create image",
        "make image of",
        "make image",
        "draw",
        "photo of",
        "picture of"
    ]

    for word in remove_words:

        prompt = prompt.replace(word, "")

    return prompt.strip()

# =========================================================
# DETECT STYLE
# =========================================================

def detect_style(prompt):

    prompt = prompt.lower()

    if "anime" in prompt:

        return """
        anime masterpiece,
        studio ghibli style,
        anime cinematic scene,
        vibrant anime colors,
        detailed anime art,
        """

    elif "cyberpunk" in prompt:

        return """
        cyberpunk,
        futuristic neon city,
        blade runner style,
        neon lights,
        sci-fi atmosphere,
        """

    elif "realistic" in prompt:

        return """
        ultra realistic,
        DSLR photography,
        photorealistic,
        cinematic lighting,
        realistic skin texture,
        """

    elif "fantasy" in prompt:

        return """
        fantasy masterpiece,
        magical environment,
        epic fantasy art,
        """

    elif "3d" in prompt:

        return """
        3D render,
        octane render,
        unreal engine,
        cinematic render,
        hyper detailed,
        """

    return """
    ultra realistic,
    cinematic lighting,
    masterpiece,
    8k,
    highly detailed,
    photorealistic,
    """

# =========================================================
# ENHANCE PROMPT
# =========================================================

def enhance_prompt(user_prompt):

    cleaned_prompt = clean_image_prompt(
        user_prompt
    )

    style = detect_style(cleaned_prompt)

    final_prompt = f"""
    {style}

    MASTERPIECE,
    BEST QUALITY,
    ULTRA DETAILED,
    8K,
    ULTRA REALISTIC,
    CINEMATIC LIGHTING,
    DRAMATIC SHADOWS,
    DEPTH OF FIELD,
    HDR,
    PROFESSIONAL PHOTOGRAPHY,
    SHARP FOCUS,
    MOVIE QUALITY,
    REALISTIC TEXTURES,
    HIGH DETAIL ENVIRONMENT,
    VISUALLY STUNNING,

    NEGATIVE:
    blurry,
    low quality,
    deformed face,
    watermark,
    text,
    ugly,

    SUBJECT:
    {cleaned_prompt}
    """

    return final_prompt

# =========================================================
# AI IMAGE GENERATION
# =========================================================

def generate_ai_image(prompt):

    try:

        seed = random.randint(
            1000,
            999999
        )

        enhanced_prompt = enhance_prompt(
            prompt
        )

        encoded_prompt = urllib.parse.quote(
            enhanced_prompt
        )

        image_url = (

            "https://image.pollinations.ai/prompt/"
            + encoded_prompt
            + f"?width=1024"
            + f"&height=1024"
            + f"&seed={seed}"
            + "&model=flux"
            + "&enhance=true"
            + "&nologo=true"
        )

        response = requests.get(

            image_url,

            timeout=120,

            headers={
                "User-Agent": "DIVU-AI"
            }
        )

        if response.status_code != 200:

            return None

        filename = (
            f"divu_{int(time.time())}_{seed}.jpg"
        )

        output_path = os.path.join(
            OUTPUT_FOLDER,
            filename
        )

        with open(output_path, "wb") as f:

            f.write(response.content)

        if os.path.getsize(output_path) < 5000:

            os.remove(output_path)

            return None

        return (
            f"http://127.0.0.1:5000/static/op/{filename}"
        )

    except Exception as e:

        print("IMAGE ERROR:", e)

        return None

# =========================================================
# IMAGE REQUEST DETECTOR
# =========================================================

def is_image_request(message):

    message = message.lower().strip()

    keywords = [

        "generate image",
        "create image",
        "make image",
        "draw",
        "anime",
        "wallpaper",
        "poster",
        "portrait",
        "realistic",
        "cinematic",
        "illustration",
        "3d render",
        "concept art",
        "photo of",
        "picture of",
        "show me",
        "make art",
        "design",
        "logo",
        "character",
        "scene",
        "fantasy",
        "cyberpunk"
    ]

    return any(
        word in message
        for word in keywords
    )

# =========================================================
# CHAT ROUTE
# =========================================================

@app.route("/chat", methods=["POST"])
def chat():

    try:

        global chat_history

        data = request.get_json()

        if not data:

            return jsonify({
                "reply": "No data received."
            })

        message = data.get(
            "message",
            ""
        ).strip()

        if not message:

            return jsonify({
                "reply": "Please enter message."
            })

        # IMAGE GENERATION

        if is_image_request(message):

            image_path = generate_ai_image(
                message
            )

            if image_path:

                return jsonify({

                    "reply":
                    "✨ DIVU AI generated your cinematic image.",

                    "image":
                    image_path
                })

            else:

                return jsonify({

                    "reply":
                    "❌ Failed to generate image."
                })

        # NORMAL CHAT

        chat_history.append({

            "role": "user",

            "content": message
        })

        if len(chat_history) > MAX_HISTORY:

            chat_history = (
                chat_history[-MAX_HISTORY:]
            )

        messages = [

            {
                "role": "system",
                "content": SYSTEM_PROMPT
            }
        ]

        messages.extend(chat_history)

        response = client.chat_completion(

            messages=messages,

            max_tokens=2500,

            temperature=0.7
        )

        reply = (
            response
            .choices[0]
            .message
            .content
        )

        reply = clean_reply(reply)

        chat_history.append({

            "role": "assistant",

            "content": reply
        })

        return jsonify({

            "reply":
            reply + "\n\n✨ Powered by DIVU AI"
        })

    except Exception as e:

        return jsonify({

            "reply":
            f"❌ Server Error: {str(e)}"
        })

# =========================================================
# PDF SUMMARY
# =========================================================

@app.route("/pdf-summary", methods=["POST"])
def pdf_summary():

    try:

        if "pdf" not in request.files:

            return jsonify({
                "reply": "No file uploaded"
            })

        file = request.files["pdf"]

        filename = secure_filename(
            file.filename
        )

        file_path = os.path.join(
            UPLOAD_FOLDER,
            filename
        )

        file.save(file_path)

        extracted_text = ""

        extension = (
            filename.split(".")[-1]
            .lower()
        )

        if extension == "pdf":

            reader = PdfReader(file_path)

            for page in reader.pages:

                try:

                    extracted_text += (
                        page.extract_text()
                        + "\n"
                    )

                except:

                    pass

        elif extension == "docx":

            doc = Document(file_path)

            for para in doc.paragraphs:

                extracted_text += (
                    para.text + "\n"
                )

        elif extension == "pptx":

            prs = Presentation(file_path)

            for slide in prs.slides:

                for shape in slide.shapes:

                    if hasattr(shape, "text"):

                        extracted_text += (
                            shape.text + "\n"
                        )

        elif extension == "txt":

            with open(

                file_path,

                "r",

                encoding="utf-8",

                errors="ignore"

            ) as f:

                extracted_text = f.read()

        else:

            return jsonify({
                "reply":
                "Unsupported file format"
            })

        extracted_text = (
            extracted_text[:7000]
        )

        response = client.chat_completion(

            messages=[

                {
                    "role": "system",

                    "content":
"""
You are DIVU AI File Analyzer.

Analyze professionally.

Provide:
- Summary
- Important Points
- Key Concepts
- Beginner Friendly Explanation
- Technical Insights
"""
                },

                {
                    "role": "user",

                    "content":
                    f"Analyze this:\n\n{extracted_text}"
                }
            ],

            max_tokens=2500,

            temperature=0.5
        )

        reply = (
            response
            .choices[0]
            .message
            .content
        )

        reply = clean_reply(reply)

        return jsonify({

            "reply":
            reply + "\n\n✨ Analyzed by DIVU AI"
        })

    except Exception as e:

        return jsonify({

            "reply":
            f"❌ File Error: {str(e)}"
        })

# =========================================================
# REAL AI IMAGE EDITING
# REPLICATE VERSION
# =========================================================

@app.route("/edit-image", methods=["POST"])
def edit_image():

    try:

        # CHECK IMAGE

        if "image" not in request.files:

            return jsonify({
                "reply":
                "No image uploaded"
            })

        image = request.files["image"]

        prompt = request.form.get(
            "prompt",
            ""
        ).strip()

        if not prompt:

            return jsonify({
                "reply":
                "Please enter prompt"
            })

        # SAVE IMAGE

        filename = secure_filename(
            image.filename
        )

        filename = (
            f"{int(time.time())}_{filename}"
        )

        original_path = os.path.join(
            UPLOAD_FOLDER,
            filename
        )

        image.save(original_path)

        # REPLICATE TOKEN

        os.environ[
            "REPLICATE_API_TOKEN"
        ] = REPLICATE_API_TOKEN

        # ADVANCED PROMPT

        enhanced_prompt = f"""

        Preserve the SAME person.
        Preserve SAME face.
        Preserve hairstyle.
        Preserve body.
        Preserve pose.
        Preserve composition.

        USER REQUEST:
        {prompt}

        Ultra realistic.
        Cinematic lighting.
        Movie quality.
        Highly detailed.
        """

        # RUN MODEL

        output = replicate.run(

            "black-forest-labs/flux-kontext-pro",

            input={

                "input_image": open(
                    original_path,
                    "rb"
                ),

                "prompt":
                enhanced_prompt
            }
        )

        # HANDLE OUTPUT

        if isinstance(output, list):

            output_url = output[0]

        else:

            output_url = output

        # DOWNLOAD IMAGE

        output_filename = (
            f"edited_{int(time.time())}.png"
        )

        output_path = os.path.join(
            OUTPUT_FOLDER,
            output_filename
        )

        img_data = requests.get(
            output_url
        ).content

        with open(output_path, "wb") as f:

            f.write(img_data)

        # SUCCESS

        return jsonify({

            "reply":
            "✨ DIVU AI edited your image successfully.",

            "image":
            f"http://127.0.0.1:5000/static/op/{output_filename}"
        })

    except Exception as e:

        return jsonify({

            "reply":
            f"❌ Error: {str(e)}"
        })

# =========================================================
# RUN
# =========================================================

if __name__ == "__main__":

    app.run(

        host="0.0.0.0",

        port=5000,

        debug=True
    )
